import sys
from PyQt5.QtWidgets import (QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
                             QListWidget, QPushButton, QFontComboBox, 
                             QSlider, QColorDialog, QTextEdit, QGraphicsView, 
                             QGraphicsScene, QGraphicsTextItem, QFileDialog)
from PyQt5.QtGui import QFont, QColor, QPixmap, QBrush
from PyQt5.QtCore import Qt, QPointF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


class ScriptureSlidesApp(QWidget):
    def __init__(self):
        super().__init__()
        self.bg_image_path = None
        self.slides_content = {}  # To store the content for each slide
        self.current_slide = 0  # Track the selected slide

        # Main layout
        self.main_layout = QHBoxLayout(self)

        # Slide List on the left
        self.slide_list = QListWidget()
        self.slide_list.addItems([f"Slide {i}" for i in range(1, 7)])  # Example slides
        self.slide_list.currentRowChanged.connect(self.on_slide_selected)  # Load slide on selection change
        self.main_layout.addWidget(self.slide_list)

        # Right side - Slide Preview and Editor
        self.right_layout = QVBoxLayout()
        
        # Font controls
        self.font_controls = self.setup_font_controls()
        self.right_layout.addLayout(self.font_controls)

        # Slide Preview (QGraphicsView)
        self.slide_preview = QGraphicsView(self)
        self.scene = QGraphicsScene(self)
        self.slide_preview.setScene(self.scene)
        self.right_layout.addWidget(self.slide_preview)

        # Add initial text to the scene (draggable text)
        self.text_item = QGraphicsTextItem("MATTHEW 28:18-20")
        self.text_item.setFont(QFont("Arial", 50))
        self.text_item.setFlag(QGraphicsTextItem.ItemIsMovable)  # Enable dragging
        self.scene.addItem(self.text_item)

        # Text Edit for Bible verses or title changes
        self.text_edit = QTextEdit(self)
        self.text_edit.setText("MATTHEW 28:18-20")
        self.text_edit.textChanged.connect(self.update_slide_text)  # Sync text with preview
        self.right_layout.addWidget(self.text_edit)

        self.main_layout.addLayout(self.right_layout)

        # Add buttons below the list
        self.add_button_layout()

    def setup_font_controls(self):
        layout = QHBoxLayout()

        # Font Combo Box
        self.font_combo = QFontComboBox()
        self.font_combo.currentFontChanged.connect(self.change_font)
        layout.addWidget(self.font_combo)

        # Font size slider
        self.font_size_slider = QSlider(Qt.Horizontal)
        self.font_size_slider.setMinimum(10)
        self.font_size_slider.setMaximum(100)
        self.font_size_slider.setValue(50)
        self.font_size_slider.valueChanged.connect(self.change_font_size)
        layout.addWidget(self.font_size_slider)

        # Bold button
        self.bold_button = QPushButton("B")
        self.bold_button.setCheckable(True)
        self.bold_button.clicked.connect(self.toggle_bold)
        layout.addWidget(self.bold_button)

        # Color picker
        self.color_picker_button = QPushButton("Color")
        self.color_picker_button.clicked.connect(self.pick_color)
        layout.addWidget(self.color_picker_button)

        return layout

    def add_button_layout(self):
        add_slide_button = QPushButton("+ Slides")
        add_slide_button.clicked.connect(self.add_slide)
        self.main_layout.addWidget(add_slide_button)

        change_bg_button = QPushButton("+ Background Image")
        change_bg_button.clicked.connect(self.change_background)
        self.main_layout.addWidget(change_bg_button)

        # Create Presentation Button
        self.create_ppt_button = QPushButton("Create Presentation", self)
        self.create_ppt_button.clicked.connect(self.create_presentation)
        self.main_layout.addWidget(self.create_ppt_button)

    def on_slide_selected(self, current_row):
        """Handles switching between slides when user selects a different slide."""
        # Save the current slide's content before switching
        if self.current_slide in self.slides_content:
            self.slides_content[self.current_slide]['text'] = self.text_edit.toPlainText()
            self.slides_content[self.current_slide]['font_size'] = self.font_size_slider.value()
            self.slides_content[self.current_slide]['bold'] = self.bold_button.isChecked()
            self.slides_content[self.current_slide]['color'] = self.text_item.defaultTextColor()

        # Switch to the newly selected slide
        self.current_slide = current_row

        # Update the text and formatting in the editor for the selected slide
        if self.current_slide in self.slides_content:
            self.text_edit.setText(self.slides_content[self.current_slide]['text'])
            self.text_item.setPlainText(self.slides_content[self.current_slide]['text'])
            self.font_size_slider.setValue(self.slides_content[self.current_slide]['font_size'])
            self.bold_button.setChecked(self.slides_content[self.current_slide]['bold'])
            color = self.slides_content[self.current_slide]['color']
            self.text_item.setDefaultTextColor(color)
        else:
            # If no content exists for the slide, clear the editor
            self.text_edit.setText(f"Slide {self.current_slide + 1}")
            self.text_item.setPlainText(f"Slide {self.current_slide + 1}")
            self.font_size_slider.setValue(50)  # Set default font size
            self.bold_button.setChecked(False)

    def update_slide_text(self):
        # Get the text from the QTextEdit and set it to the QGraphicsTextItem
        self.text_item.setPlainText(self.text_edit.toPlainText())

    def change_font(self, font):
        self.text_item.setFont(font)

    def change_font_size(self):
        size = self.font_size_slider.value()
        font = self.text_item.font()
        font.setPointSize(size)
        self.text_item.setFont(font)

    def toggle_bold(self):
        font = self.text_item.font()
        font.setBold(self.bold_button.isChecked())
        self.text_item.setFont(font)

    def pick_color(self):
        color = QColorDialog.getColor()
        if color.isValid():
            self.text_item.setDefaultTextColor(color)

    def add_slide(self):
        new_slide_number = self.slide_list.count() + 1
        self.slide_list.addItem(f"Slide {new_slide_number}")
        self.slides_content[new_slide_number - 1] = {'text': f"Slide {new_slide_number}", 
                                                      'font_size': 50, 
                                                      'bold': False, 
                                                      'color': QColor(Qt.black)}

    def change_background(self):
        # Use file dialog to select an image and set it as background
        self.bg_image_path, _ = QFileDialog.getOpenFileName(self, 'Select Background Image', '', 'Images (*.png *.jpg *.bmp)')
        if self.bg_image_path:
            bg_pixmap = QPixmap(self.bg_image_path)
            bg_brush = QBrush(bg_pixmap)  # Convert QPixmap to QBrush
            self.scene.setBackgroundBrush(bg_brush)  # Set background image

    def create_presentation(self):
        # Create a new presentation
        prs = Presentation()

        # Loop through each slide and add content
        for i in range(self.slide_list.count()):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add background image if it was selected
            if self.bg_image_path:
                slide.shapes.add_picture(self.bg_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # Add text from stored slide content
            if i in self.slides_content:
                slide_text = self.slides_content[i]['text']
                font_size = self.slides_content[i]['font_size']
                is_bold = self.slides_content[i]['bold']
                color = self.slides_content[i]['color']

                # Add text to slide
                left = Inches(1)
                top = Inches(1)
                width = Inches(8)
                height = Inches(2)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                p = text_frame.add_paragraph()
                p.text = slide_text

                # Apply font styles
                p.font.size = Pt(font_size)
                p.font.bold = is_bold
                p.font.color.rgb = RGBColor(color.red(), color.green(), color.blue())

        # Save the presentation
        save_path, _ = QFileDialog.getSaveFileName(self, 'Save Presentation', '', 'PowerPoint Files (*.pptx)')
        if save_path:
            prs.save(save_path)


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = ScriptureSlidesApp()
    window.setWindowTitle("Scripture Slides.exe")
    window.resize(1200, 700)
    window.show()
    sys.exit(app.exec_())
