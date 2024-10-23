import os
import sys
import requests
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QPushButton, QLabel, QLineEdit, QFileDialog, 
    QComboBox, QColorDialog, QTextEdit, QFontComboBox, QMainWindow, QAction, 
    QMessageBox, QHBoxLayout, QListWidget, QListWidgetItem, QInputDialog, QStackedWidget, 
)
from PyQt5.QtGui import QPixmap, QPainter, QColor, QFont
from PyQt5.QtCore import Qt
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

class DraggableLabel(QLabel):
    def __init__(self, text, parent=None):
        super().__init__(text, parent)
        self.setStyleSheet("background-color: rgba(255, 255, 255, 0);")  # Transparent background
        self.setAlignment(Qt.AlignCenter)
        self.setWordWrap(True)
        self._drag_active = False
        self._drag_position = None

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._drag_active = True
            self._drag_position = event.globalPos() - self.frameGeometry().topLeft()
            event.accept()

    def mouseMoveEvent(self, event):
        if event.buttons() == Qt.LeftButton and self._drag_active:
            self.move(event.globalPos() - self._drag_position)
            event.accept()

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._drag_active = False
            event.accept()

class SlidePreviewWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.background = None
        layout = QVBoxLayout(self)
        self.setLayout(layout)

    def set_background(self, background_path):
        """Sets the background image."""
        self.background = QPixmap(background_path)
        self.update()  # Repaint widget

    def paintEvent(self, event):
        """Override to draw the background."""
        painter = QPainter(self)
        if self.background:
            scaled_bg = self.background.scaled(self.size(), Qt.KeepAspectRatioByExpanding)
            painter.drawPixmap(self.rect(), scaled_bg)

        # Call parent class paintEvent
        super().paintEvent(event)


class ScriptureSlidesApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.api_token = os.getenv("ESV_API_TOKEN")
        self.books_of_the_bible = [
            "Genesis", "Exodus", "Leviticus", "Numbers", "Deuteronomy", "Joshua", "Judges", "Ruth", "1 Samuel",
            "2 Samuel", "1 Kings", "2 Kings", "1 Chronicles", "2 Chronicles", "Ezra", "Nehemiah", "Esther", "Job",
            "Psalms", "Proverbs", "Ecclesiastes", "Song of Solomon", "Isaiah", "Jeremiah", "Lamentations", "Ezekiel",
            "Daniel", "Hosea", "Joel", "Amos", "Obadiah", "Jonah", "Micah", "Nahum", "Habakkuk", "Zephaniah",
            "Haggai", "Zechariah", "Malachi", "Matthew", "Mark", "Luke", "John", "Acts", "Romans", "1 Corinthians",
            "2 Corinthians", "Galatians", "Ephesians", "Philippians", "Colossians", "1 Thessalonians", "2 Thessalonians",
            "1 Timothy", "2 Timothy", "Titus", "Philemon", "Hebrews", "James", "1 Peter", "2 Peter", "1 John", "2 John",
            "3 John", "Jude", "Revelation"
        ]
        self.verses_list = []  # List to store multiple verse configurations
        self.slides_data = []  # To hold slides with all configurations
        self.current_slide_index = None  # Track the currently edited slide
        self.initUI()

    def initUI(self):
        self.setWindowTitle("Scripture Slides")

        # Set the main widget
        main_widget = QWidget(self)
        main_layout = QVBoxLayout(main_widget)

        # Menu bar
        menubar = self.menuBar()

        # File menu (Presets, Help, Version)
        fileMenu = menubar.addMenu('File')

        # Presets action
        presets_action = QAction('Presets', self)
        presets_action.triggered.connect(self.show_presets)
        fileMenu.addAction(presets_action)

        # Help action
        help_action = QAction('Help', self)
        help_action.triggered.connect(self.show_help)
        fileMenu.addAction(help_action)

        # Version action
        version_action = QAction('Version', self)
        version_action.triggered.connect(self.show_version)
        fileMenu.addAction(version_action)

        # Bible search section
        search_layout = QHBoxLayout()

        # Create a combo box for book selection
        self.book_input = QComboBox(self)
        self.book_input.addItems(self.books_of_the_bible)
        self.book_input.currentTextChanged.connect(self.update_chapters)
        search_layout.addWidget(self.book_input)

        # Create a combo box for chapter selection
        self.chapter_input = QComboBox(self)
        self.chapter_input.currentTextChanged.connect(self.update_verses)
        search_layout.addWidget(self.chapter_input)

        # Create a combo box for verse selection
        self.verse_input = QComboBox(self)
        search_layout.addWidget(self.verse_input)

        # Button to add verse to the list
        self.add_verse_button = QPushButton("Add Verse", self)
        self.add_verse_button.clicked.connect(self.add_verse_to_list)
        search_layout.addWidget(self.add_verse_button)

        main_layout.addLayout(search_layout)

        # Create a horizontal layout for the slide preview and editor
        preview_layout = QHBoxLayout()

        # List to display selected verses
        self.verses_list_widget = QListWidget(self)
        self.verses_list_widget.itemClicked.connect(self.preview_slide)  # On-click to preview slide
        preview_layout.addWidget(self.verses_list_widget)

        # Slide preview and edit section
        self.slide_preview_stack = QStackedWidget(self)
        preview_layout.addWidget(self.slide_preview_stack)

        # Add the preview layout to the main layout
        main_layout.addLayout(preview_layout)

        # Button to remove a selected verse from the list
        self.remove_verse_button = QPushButton("Remove Slide", self)
        self.remove_verse_button.clicked.connect(self.remove_selected_verse)
        main_layout.addWidget(self.remove_verse_button)

        # Button to save slide edits
        self.save_slide_button = QPushButton('Save Slide Changes', self)
        self.save_slide_button.clicked.connect(self.save_slide_changes)
        main_layout.addWidget(self.save_slide_button)

        # Editor section: font family and size selection for title and verse
        self.font_family_title = QFontComboBox(self)
        main_layout.addWidget(self.font_family_title)

        self.font_size_title = QComboBox(self)
        self.font_size_title.addItems([str(i) for i in range(24, 73, 2)])  # Sizes 24 to 72
        main_layout.addWidget(self.font_size_title)

        self.font_family_verse = QFontComboBox(self)
        main_layout.addWidget(self.font_family_verse)

        self.font_size_verse = QComboBox(self)
        self.font_size_verse.addItems([str(i) for i in range(16, 49, 2)])  # Sizes 16 to 48
        main_layout.addWidget(self.font_size_verse)

        # Button to change title color
        self.title_color_button = QPushButton('Change Title Color', self)
        self.title_color_button.clicked.connect(self.pick_title_color)
        main_layout.addWidget(self.title_color_button)

        # Button to change verse color
        self.verse_color_button = QPushButton('Change Verse Color', self)
        self.verse_color_button.clicked.connect(self.pick_verse_color)
        main_layout.addWidget(self.verse_color_button)

        # Label to display selected background
        self.label = QLabel("No backgrounds selected", self)
        main_layout.addWidget(self.label)

        # Button to upload background images
        self.bg_button = QPushButton('Upload Background Images', self)
        self.bg_button.clicked.connect(self.upload_backgrounds)
        main_layout.addWidget(self.bg_button)

        # Button to create an empty slide with background
        self.create_empty_slide_button = QPushButton('Create Empty Slide', self)
        self.create_empty_slide_button.clicked.connect(self.create_empty_slide)
        main_layout.addWidget(self.create_empty_slide_button)

        # LineEdit for specifying the file name
        self.file_name_input = QLineEdit(self)
        self.file_name_input.setPlaceholderText("Enter file name for the PowerPoint (without extension)")
        main_layout.addWidget(self.file_name_input)

        # Button to create the PowerPoint
        self.ppt_button = QPushButton('Create PowerPoint', self)
        self.ppt_button.clicked.connect(self.create_presentation)
        main_layout.addWidget(self.ppt_button)

        self.setCentralWidget(main_widget)

        # Color placeholders
        self.title_color = None
        self.verse_color = None

        # List to hold multiple background images
        self.background_images = []

        # Initialize chapters and verses
        self.update_chapters()

    def create_empty_slide(self):
        """Creates a slide with just the background image."""
        if not self.background_images:
            QMessageBox.warning(self, 'Error', 'Please upload a background image first.')
            return

        # Add an empty slide configuration with background only
        self.slides_data.append({
            "book": None,
            "chapter": None,
            "verse": None,
            "title_font": None,
            "title_size": None,
            "verse_font": None,
            "verse_size": None,
            "title_color": None,
            "verse_color": None,
            "background": self.background_images[0]  # Choose the first background
        })

        self.verses_list_widget.addItem('Empty Slide')
        QMessageBox.information(self, 'Slide Created', 'Empty slide with background created successfully.')

    def update_chapters(self):
        """Update the chapter combo box based on the selected book."""
        selected_book = self.book_input.currentText()
        self.chapter_input.clear()
        self.chapter_input.addItems([str(i) for i in range(1, 51)])  # Assume 50 chapters per book
        self.update_verses()

    def update_verses(self):
        """Update the verse combo box based on the selected chapter."""
        self.verse_input.clear()
        self.verse_input.addItems([str(i) for i in range(1, 31)])

    def add_verse_to_list(self):
        """Add the current verse selection to the list and prepare for editing."""
        book = self.book_input.currentText()
        chapter = self.chapter_input.currentText()
        verse = self.verse_input.currentText()

        verse_config = f"{book} {chapter}:{verse}"
        verse_text = self.fetch_bible_verse(book, chapter, verse)

        # Save the verse data
        self.slides_data.append({
            "book": book,
            "chapter": chapter,
            "verse": verse,
            "verse_text": verse_text,
            "title_font": None,
            "title_size": None,
            "verse_font": None,
            "verse_size": None,
            "title_color": None,
            "verse_color": None,
            "background": self.background_images[0] if self.background_images else None
        })

        self.verses_list_widget.addItem(verse_config)

    def fetch_bible_verse(self, book, chapter, verse):
        """Fetch the Bible verse text using the ESV API."""
        if not self.api_token:
            QMessageBox.warning(self, 'Error', 'No API token found. Please set your ESV API token.')
            return None

        # Construct the API URL
        url = f"https://api.esv.org/v3/passage/text/?q={book}%20{chapter}:{verse}"
        headers = {
            "Authorization": f"Token {self.api_token}"
        }
        try:
            response = requests.get(url, headers=headers)
            if response.status_code == 200:
                data = response.json()
                passage_text = data["passages"][0]
                return passage_text.strip()
            else:
                QMessageBox.warning(self, 'Error', 'Failed to fetch verse. Check API token and internet connection.')
        except requests.RequestException as e:
            QMessageBox.warning(self, 'Error', f"An error occurred: {str(e)}")
        return None

    def preview_slide(self, item):
        """Preview the selected slide based on the item clicked in the list."""
        selected_index = self.verses_list_widget.row(item)
        self.current_slide_index = selected_index
        slide_data = self.slides_data[selected_index]

        # Remove all widgets from QStackedWidget by replacing them with new empty widgets
        while self.slide_preview_stack.count() > 0:
            widget = self.slide_preview_stack.widget(0)
            self.slide_preview_stack.removeWidget(widget)
            widget.deleteLater()

        # Load and display the background image in a QLabel
        if slide_data.get("background"):
            background_image_path = slide_data["background"]
            pixmap = QPixmap(background_image_path)
            preview_label = QLabel()
            preview_label.setPixmap(pixmap.scaled(
                self.slide_preview_stack.size(),  # Use size method for current dimensions
                Qt.KeepAspectRatio  # Use Qt.KeepAspectRatio for proper scaling
            ))
            self.slide_preview_stack.addWidget(preview_label)

        # Create a QLabel to display the verse and title
        if slide_data.get("book"):
            verse_info = f"{slide_data['book']} {slide_data['chapter']}:{slide_data['verse']}"
            verse_text = slide_data["verse_text"]
        else:
            verse_info = "Empty Slide"
            verse_text = ""

        # Create the title label with safety checks
        title_font = slide_data.get("title_font", "Arial")  # Default to Arial if None
        title_size = slide_data.get("title_size", 24)  # Default size if None

        # Ensure title_size is a valid integer
        if title_size is None:
            title_size = 24

        # Create the title label with proper HTML formatting
        title_label = QLabel(f"<b>{verse_info}</b><br>{verse_text}", self)
        title_label.setWordWrap(True)
        title_label.setAlignment(Qt.AlignCenter)  # Align text in the center
        title_label.setFont(QFont(title_font, title_size))  # Create font safely

        # Add title label to the preview stack
        self.slide_preview_stack.addWidget(title_label)

        # Set current index to the newly added preview
        self.slide_preview_stack.setCurrentIndex(0)




    def save_slide_changes(self):
        """Save changes made to the currently selected slide."""
        if self.current_slide_index is None:
            QMessageBox.warning(self, 'Error', 'No slide selected to edit.')
            return

        slide_data = self.slides_data[self.current_slide_index]

        # Save title font and size
        slide_data["title_font"] = self.font_family_title.currentFont().family()
        slide_data["title_size"] = int(self.font_size_title.currentText())

        # Save verse font and size
        slide_data["verse_font"] = self.font_family_verse.currentFont().family()
        slide_data["verse_size"] = int(self.font_size_verse.currentText())

        # Save title and verse colors
        slide_data["title_color"] = self.title_color
        slide_data["verse_color"] = self.verse_color

        QMessageBox.information(self, 'Success', 'Slide changes saved successfully.')

    def pick_title_color(self):
        """Pick a color for the title text."""
        color = QColorDialog.getColor()
        if color.isValid():
            self.title_color = color
            self.title_color_button.setStyleSheet(f"background-color: {color.name()};")

    def pick_verse_color(self):
        """Pick a color for the verse text."""
        color = QColorDialog.getColor()
        if color.isValid():
            self.verse_color = color
            self.verse_color_button.setStyleSheet(f"background-color: {color.name()};")

    def upload_backgrounds(self):
        """Upload background images for the slides."""
        options = QFileDialog.Options()
        files, _ = QFileDialog.getOpenFileNames(self, "Select Background Images", "", 
                                                "Images (*.png *.jpg *.jpeg *.bmp)", options=options)
        if files:
            self.background_images = files
            self.label.setText(f"{len(files)} background(s) selected.")
            QMessageBox.information(self, 'Success', 'Background images uploaded successfully.')

    def remove_selected_verse(self):
        """Remove the selected verse from the list and update the slide data."""
        selected_item = self.verses_list_widget.currentRow()
        if selected_item != -1:
            self.verses_list_widget.takeItem(selected_item)
            self.slides_data.pop(selected_item)

            QMessageBox.information(self, 'Success', 'Selected verse removed successfully.')
        else:
            QMessageBox.warning(self, 'Error', 'No verse selected to remove.')

    def show_presets(self):
        """Show available presets for slide design."""
        QMessageBox.information(self, 'Presets', 'No presets available yet.')

    def show_help(self):
        """Display help information."""
        QMessageBox.information(self, 'Help', 'Instructions on how to use the Scripture Slides app...')

    def show_version(self):
        """Display the version information."""
        QMessageBox.information(self, 'Version', 'Scripture Slides v1.0')

    def create_presentation(self):
        """Create a PowerPoint presentation with the configured slides."""
        if not self.file_name_input.text():
            QMessageBox.warning(self, 'Error', 'Please enter a file name for the presentation.')
            return

        # Create a new PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions to 1920x1080 (20 inches wide, 11.25 inches high)
        prs.slide_width = Inches(20)
        prs.slide_height = Inches(11.25)

        for slide_data in self.slides_data:
            slide_layout = prs.slide_layouts[5]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)

            # Add background image if available
            if slide_data.get("background"):
                img_path = slide_data["background"]
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # Add title (book and chapter info) if available
            if slide_data.get("book"):
                title_placeholder = slide.shapes.title
                title_placeholder.text = f"{slide_data['book']} {slide_data['chapter']}:{slide_data['verse']}"
                
                # Get title size, default to 24 if None
                title_size = slide_data.get("title_size")
                if title_size is None:
                    title_size = 24  # Default value if None
                title_placeholder.text_frame.paragraphs[0].font.size = Pt(title_size)
                
                # Set title font if specified
                if slide_data.get("title_font"):
                    title_placeholder.text_frame.paragraphs[0].font.name = slide_data["title_font"]
                
                # Set title color if specified
                if slide_data.get("title_color"):
                    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        slide_data["title_color"].red(),
                        slide_data["title_color"].green(),
                        slide_data["title_color"].blue()
                    )

            # Add verse text if available
            if slide_data.get("verse_text"):
                left = Inches(0.5)
                top = Inches(2.0)
                width = Inches(8.0)
                height = Inches(4.0)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.text = slide_data["verse_text"]
                
                # Get verse size, default to 18 if None
                verse_size = slide_data.get("verse_size")
                if verse_size is None:
                    verse_size = 18  # Default value if None
                text_frame.paragraphs[0].font.size = Pt(verse_size)
                
                # Set verse font if specified
                if slide_data.get("verse_font"):
                    text_frame.paragraphs[0].font.name = slide_data["verse_font"]
                
                # Set verse color if specified
                if slide_data.get("verse_color"):
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        slide_data["verse_color"].red(),
                        slide_data["verse_color"].green(),
                        slide_data["verse_color"].blue()
                    )

        # Save the presentation
        file_name = self.file_name_input.text() + ".pptx"
        prs.save(file_name)

        QMessageBox.information(self, 'Success', f'Presentation saved as {file_name}.')


if __name__ == '__main__':
    app = QApplication(sys.argv)
    window = ScriptureSlidesApp()
    window.show()
    sys.exit(app.exec_())
