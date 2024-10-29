import sys
import os
import time
import warnings
from PyQt5.QtWidgets import (QMainWindow, QApplication, QFileDialog, QFontComboBox, QPushButton,
                             QListWidget, QColorDialog, QGraphicsScene, QGraphicsTextItem, QGraphicsView, QTextEdit, QMenu)
from PyQt5.QtGui import QFont, QColor, QPixmap, QBrush, QIcon, QTextCursor, QTextCharFormat
from PyQt5.uic import loadUi
from PyQt5.QtCore import Qt, pyqtSignal
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PyQt5.QtWidgets import QGraphicsPixmapItem, QDialog, QLabel, QVBoxLayout, QPushButton, QSpinBox # Added for DraggableImageItem
import win32com.client as win32
import msvcrt

warnings.simplefilter("ignore", DeprecationWarning)

from PIL import Image, ImageDraw, ImageFont

class DraggableTextItem(QGraphicsTextItem):
    def __init__(self, text):
        super().__init__(text)
        self.setTextInteractionFlags(Qt.TextEditorInteraction)  # Allow direct editing
        self.setFlag(QGraphicsTextItem.ItemIsMovable, True)      # Enable dragging
        self.setFlag(QGraphicsTextItem.ItemIsSelectable, True)   # Enable selection

class DraggableImageItem(QGraphicsPixmapItem):  # QGraphicsPixmapItem imported here
    def __init__(self, pixmap):
        super().__init__(pixmap)
        self.setFlags(QGraphicsPixmapItem.ItemIsMovable | QGraphicsPixmapItem.ItemIsSelectable)

class VerseRepeatWindow(QDialog):
    def __init__(self, parent=None):
        super(VerseRepeatWindow, self).__init__(parent)
        self.setWindowTitle("Add Verse")

        # Load the UI file for the VerseRepeatWindow (optional if you don't have one)
        loadUi("Scripture-Slides/AddVerse.ui", self)

class SlideShowWindow(QMainWindow):
    closed = pyqtSignal()

    def __init__(self, parent=None):
        super(SlideShowWindow, self).__init__(parent)
        loadUi("Scripture-Slides/SlideShowWindow.ui", self)
        
        self.backButton.clicked.connect(self.close)

        # Initialize QGraphicsScene and other attributes needed for live edit
        self.scene = QGraphicsScene(self.graphicsView)
        self.graphicsView.setScene(self.scene)
        
        # Track currently selected slide and previews
        self.current_slide_index = None
        self.slide_previews = {}

        # Connect buttons for live editing
        self.LiveEditBtn.clicked.connect(self.enter_live_edit_mode)
        self.LiveBtn.clicked.connect(self.apply_live_changes)

        # New text box for copying/pasting text
        self.textInputBox = QTextEdit(self)
        self.textInputBox.setPlaceholderText("Copy and paste text here...")
        self.textInputBox.setFixedHeight(100)
        layout = QVBoxLayout(self.graphicsView)
        layout.addWidget(self.textInputBox)
        self.setLayout(layout)
    
    def enter_live_edit_mode(self):
        """Enable live editing of the selected slide."""
        selected_items = self.slideListWidget.selectedItems()
        if selected_items:
            selected_item = selected_items[0].text()
            print(f"Entering live edit mode for {selected_item}")

            # Enable text and image items to be moved and edited
            for item in self.scene.items():
                if isinstance(item, (DraggableTextItem, DraggableImageItem)):
                    item.setFlag(QGraphicsTextItem.ItemIsMovable, True)

    def apply_live_changes(self):
        """Apply live changes to the selected slide in the PowerPoint file."""
        selected_items = self.slideListWidget.selectedItems()
        if selected_items:
            selected_item = selected_items[0].text()
            selected_index = int(selected_item.split(" ")[1]) - 1
            slide = self.parent().prs.slides[selected_index]
            
            # Clear existing shapes on the slide
            for shape in slide.shapes:
                sp = shape._element
                sp.getparent().remove(sp)

            # Add copied text from the input box to slide
            pasted_text = self.textInputBox.toPlainText()
            text_item = DraggableTextItem(pasted_text)
            text_item.setFont(QFont("Arial", 20))
            self.scene.addItem(text_item)

            # Save preview and update slide
            self.parent().save_slide_preview()
            print(f"Changes applied to {selected_item}")

    def load_slide_previews(self, slide_previews):
        """Load slide previews from the main window to display."""
        self.slide_previews = slide_previews
        self.update_slide_list()
        
    def update_slide_list(self):
        """Populate the slide list with preview images."""
        self.slideListWidget.clear()
        for slide_name in self.slide_previews:
            self.slideListWidget.addItem(slide_name)
        self.slideListWidget.itemSelectionChanged.connect(self.display_slide_in_graphics_view)

    def display_slide_in_graphics_view(self):
        """Display the selected slide in QGraphicsView."""
        selected_items = self.slideListWidget.selectedItems()
        if selected_items:
            selected_item = selected_items[0].text()
            image_path = self.slide_previews.get(selected_item)
            if image_path:
                self.show_slide_preview(image_path)
            else:
                print(f"No preview available for {selected_item}")
                
    def show_slide_preview(self, image_path):
        """Display an image in the QGraphicsView."""
        pixmap = QPixmap(image_path)
        self.scene.clear()
        self.scene.addPixmap(pixmap)
        self.graphicsView.fitInView(self.scene.itemsBoundingRect(), Qt.KeepAspectRatio)
        
    def enter_live_edit_mode(self):
        """Enable live editing of the selected slide."""
        # Logic for entering live edit mode based on the selected slide
        print("Entering live edit mode.")
        # Additional code here for text/image manipulation in live edit mode
        
    def apply_live_changes(self):
        """Apply live changes to the selected slide."""
        print("Applying live changes to the slideshow.")
        # Logic to save changes to the main slideshow (update presentation object, etc.)
        
    def closeEvent(self, event):
        # Emit the closed signal
        self.closed.emit()
        super(SlideShowWindow, self).closeEvent(event)

class ScriptureSlides(QMainWindow):
    def __init__(self):
        super(ScriptureSlides, self).__init__()
        self.setWindowFlags(Qt.Window | Qt.WindowMinimizeButtonHint | Qt.WindowCloseButtonHint)
        loadUi("Scripture-Slides/ScriptureSlides.ui", self)

        self.setStyleSheet(""" 
                QComboBox#fontComboBox {
            font-family: 'Inter 18pt', sans-serif;
            font-size: 10pt; /* Change font size to points */
            color: black;
            background-color: white;
            border: 1px solid #ccc;
            padding: 5px;
            border-radius: 8px;
        }

        QComboBox#fontComboBox::drop-down {
            background-color: transparent; /* Change drop-down background */
            width: 30px; /* Adjust drop-down width */
        }

        QComboBox#fontComboBox::drop-down:hover {
        background-color: #EFEFEF;
            border-radius: 8px;
        }

        QComboBox#fontComboBox::down-arrow {
            image: url(C:/Users/markt/Scripture-Slides/Scripture-Slides/assets/arrowDown.png); /* Path to your custom arrow image */
            width: 18px; /* Adjust width if needed */
            height: 18px; /* Adjust height if needed */
        }

        QComboBox#fontComboBox QAbstractItemView {
            font-family: 'Inter 18pt', sans-serif;
            font-size: 10pt; /* Set point size for list items */
            color: black;
            background-color: white;
            border: none;
        }

        """)

        # Initialize presentation object and QGraphicsScene only once
        self.prs = Presentation()
        self.prs.slide_width = Inches(20)
        self.prs.slide_height = Inches(11.25)
        self.slide_count = 0
        self.slide_previews = {}

        # Initialize QGraphicsScene for preview
        self.scene = QGraphicsScene(self.graphicsView)
        self.graphicsView.setScene(self.scene)

        # Connect buttons to functions
        self.addSlideBtn.clicked.connect(self.add_slide)
        self.addBackgroundImageBtn.clicked.connect(self.add_background_image)
        self.createPresentationBtn.clicked.connect(self.create_presentation)
        self.slideListWidget.itemSelectionChanged.connect(self.display_slide_in_graphics_view)
        self.addTextBtn.clicked.connect(self.add_text_item) 
        self.slideListWidget.setContextMenuPolicy(Qt.CustomContextMenu)
        self.slideListWidget.customContextMenuRequested.connect(self.open_context_menu)
        self.current_slide = None
        self.VerseRepeatBtn.clicked.connect(self.open_verse_repeat_window)
        self.verse_repeat_window = None
        self.SlideShowBtn.clicked.connect(self.open_slideshow_window)
        self.slideshow_window = None
        self.addTextBtn.clicked.connect(self.add_text_item) 
        self.fontComboBox.currentFontChanged.connect(self.change_font_family)
        self.DecreaseFontSize.clicked.connect(self.decrease_font_size)
        self.IncreaseFontSize.clicked.connect(self.increase_font_size)
        self.BoldBtn.clicked.connect(self.toggle_bold)
        self.colorWheel.clicked.connect(self.change_font_color)
        self.AlignLeft.clicked.connect(lambda: self.change_alignment(Qt.AlignLeft))
        self.AlignCenter.clicked.connect(lambda: self.change_alignment(Qt.AlignCenter))
        self.AlignRight.clicked.connect(lambda: self.change_alignment(Qt.AlignRight))
        self.AlignJustify.clicked.connect(lambda: self.change_alignment(Qt.AlignJustify))

        # Default font settings
        self.current_font = QFont("Arial", 20)
        self.current_color = QColor(0, 0, 0)  # Default to black
        self.current_alignment = Qt.AlignLeft


    def open_slideshow_window(self):
        """Open the Slideshow window and hide the main window."""
        if not self.slideshow_window or not self.slideshow_window.isVisible():
            self.slideshow_window = SlideShowWindow(self)
            self.slideshow_window.load_slide_previews(self.slide_previews)  # Pass slide previews here
            self.hide()  # Hide the main window
            self.slideshow_window.show()
            self.slideshow_window.closed.connect(self.show_main_window)
        else:
            self.slideshow_window.raise_()

    def show_main_window(self):
        """Show the main window when slideshow window is closed."""
        self.show()

    def close_slideshow_window(self):
        print("Slideshow Window closed.")
        self.slideshow_window = None

    def open_verse_repeat_window(self):
        """Open the Verse Repeat window."""
        # Check if window already exists to prevent multiple instances
        if self.verse_repeat_window is None or not self.verse_repeat_window.isVisible():
            self.verse_repeat_window = VerseRepeatWindow(self)
            self.verse_repeat_window.finished.connect(self.on_verse_repeat_window_closed)
            self.verse_repeat_window.show()
        else:
            self.verse_repeat_window.raise_()  # Bring the existing window to the front

    def on_verse_repeat_window_closed(self):
        """Handle the closing of the Verse Repeat window."""
        print("Verse Repeat window closed.")
        self.verse_repeat_window = None  # Reset the window instance when closed


    #TEXT FORMATTING BUTTONS!!

    def add_text_item(self):
        """Add a new draggable text item to the graphics view if it doesn't already exist."""
        if hasattr(self, 'current_text_item') and self.current_text_item in self.scene.items():
            print("Text item already added.")
            return  # Avoid adding duplicate text items

        text_item = DraggableTextItem("Editable Text")
        text_item.setFont(self.current_font)
        text_item.setDefaultTextColor(self.current_color)
        self.scene.addItem(text_item)
        text_item.setTextWidth(200)  # Optional, for wrapping text
        self.current_text_item = text_item  # Track the currently edited text item

    def apply_text_formatting(self):
        """Apply all current font settings to the tracked text item without creating duplicates."""
        if hasattr(self, 'current_text_item') and self.current_text_item:
            # Apply formatting only to the existing text item
            self.current_text_item.setFont(self.current_font)
            self.current_text_item.setDefaultTextColor(self.current_color)
            cursor = self.current_text_item.textCursor()
            cursor.select(QTextCursor.Document)
            text_format = QTextCharFormat()
            text_format.setFont(self.current_font)
            text_format.setForeground(self.current_color)
            cursor.mergeCharFormat(text_format)
        else:
            print("No text item selected for formatting.")

    def change_font_family(self, font):
        """Update the font family for the current text item only."""
        self.current_font.setFamily(font.family())
        self.apply_text_formatting()


    def change_font_family(self, font):
        """Update the font family for the selected text item."""
        self.current_font.setFamily(font.family())
        if self.current_text_item:
            self.current_text_item.setFont(self.current_font)

    def increase_font_size(self):
        """Increase the font size by 1 point."""
        self.current_font.setPointSize(self.current_font.pointSize() + 1)
        if self.current_text_item:
            self.current_text_item.setFont(self.current_font)

    def decrease_font_size(self):
        """Decrease the font size by 1 point."""
        if self.current_font.pointSize() > 1:  # Ensure size stays positive
            self.current_font.setPointSize(self.current_font.pointSize() - 1)
        if self.current_text_item:
            self.current_text_item.setFont(self.current_font)

    def change_font_color(self):
        """Open color dialog to select a font color."""
        color = QColorDialog.getColor(self.current_color, self)
        if color.isValid():
            self.current_color = color
            if self.current_text_item:
                self.current_text_item.setDefaultTextColor(self.current_color)

    def toggle_bold(self):
        """Toggle bold for the current font."""
        self.current_font.setBold(not self.current_font.bold())
        if self.current_text_item:
            self.current_text_item.setFont(self.current_font)

    def set_text_alignment(self, alignment):
        """Set text alignment for the current text item."""
        self.current_alignment = alignment
        if self.current_text_item:
            self.current_text_item.setTextAlignment(self.current_alignment)
    

    def add_slide(self):
        """Add a new blank slide and save its preview."""
        slide_layout = self.prs.slide_layouts[6]  # Layout for a blank slide
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1

        slide_item = f"Slide {self.slide_count}"
        self.slideListWidget.addItem(slide_item)
        
        self.save_slide_preview()  # Capture preview (using placeholder data)
        print(f"{slide_item} added!")

    def add_background_image(self):
        """Add a draggable background image to the QGraphicsView and PowerPoint slide."""
        selected_items = self.slideListWidget.selectedItems()
        if not selected_items:
            print("No slide selected. Add a slide first.")
            return

        selected_index = self.slideListWidget.currentRow()
        self.current_slide = self.prs.slides[selected_index]  # Get the actual slide from the selected index

        # Open file dialog to select an image
        image_path, _ = QFileDialog.getOpenFileName(self, 'Open Background Image', '', 'Image files (*.jpg *.png)')
        if image_path:
            # Add image to QGraphicsView for a preview
            pixmap = QPixmap(image_path)
            image_item = DraggableImageItem(pixmap)
            
            # Clear any previous background image to avoid stacking images
            self.scene.clear()
            
            # Add the new image as a draggable item
            self.scene.addItem(image_item)
            
            # Adjust the view to fit the new image
            self.graphicsView.fitInView(self.scene.itemsBoundingRect(), Qt.KeepAspectRatio)
            print(f"Draggable background image {image_path} added to the graphics view.")

            # Update preview with new image path
            self.save_slide_preview(image_path)

            # Set the background image in PowerPoint slide
            left, top, width, height = 0, 0, self.prs.slide_width, self.prs.slide_height
            self.current_slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"Background image {image_path} added to slide in PowerPoint presentation.")


    def save_slide_preview(self, image_path=None):
        """Simulate the current slide preview and save it as an image."""
        preview_folder = "./slide_previews/"
        os.makedirs(preview_folder, exist_ok=True)
        
        slide_number = self.slideListWidget.currentRow() + 1
        preview_image_path = f"{preview_folder}slide_{slide_number}.png"
        
        slide_width, slide_height = 1280, 720
        image = Image.new("RGB", (slide_width, slide_height), "white")
        draw = ImageDraw.Draw(image)

        if image_path:
            background = Image.open(image_path)
            background = background.resize((slide_width, slide_height))
            image.paste(background, (0, 0))

        text = f"Slide {slide_number}"
        font = ImageFont.load_default()
        draw.text((50, 50), text, fill=(0, 0, 0), font=font)

        image.save(preview_image_path)
        print(f"Preview for Slide {slide_number} saved at {preview_image_path}")
        
        self.slide_previews[f"Slide {slide_number}"] = preview_image_path

    def display_image_in_graphics_view(self, image_path):
        """Display an image in the QGraphicsView."""
        pixmap = QPixmap(image_path)
        self.scene.clear()
        self.scene.addPixmap(pixmap)
        self.graphicsView.fitInView(self.scene.itemsBoundingRect(), Qt.KeepAspectRatio)

    def display_slide_in_graphics_view(self):
        """Display the selected slide in QGraphicsView."""
        selected_items = self.slideListWidget.selectedItems()
        if selected_items:
            selected_item = selected_items[0].text()
            image_path = self.slide_previews.get(selected_item)
            if image_path:
                self.display_image_in_graphics_view(image_path)
            else:
                print(f"No preview available for {selected_item}")

    def open_context_menu(self, position):
        """Open a custom context menu for deleting slides."""
        context_menu = QMenu(self)
        delete_action = context_menu.addAction("Delete Slide")
        action = context_menu.exec_(self.slideListWidget.mapToGlobal(position))
        
        if action == delete_action:
            self.delete_slide()

    def delete_slide(self):
        """Delete the selected slide from the list and presentation."""
        selected_row = self.slideListWidget.currentRow()
        if selected_row >= 0:
            selected_item = self.slideListWidget.takeItem(selected_row)
            slide_name = selected_item.text()
            
            preview_path = self.slide_previews.pop(slide_name, None)
            if preview_path and os.path.exists(preview_path):
                os.remove(preview_path)
                print(f"Deleted preview image: {preview_path}")

            new_prs = Presentation()
            new_prs.slide_width = self.prs.slide_width
            new_prs.slide_height = self.prs.slide_height
            
            for i, slide in enumerate(self.prs.slides):
                if i != selected_row:
                    new_slide = new_prs.slides.add_slide(slide.slide_layout)
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            new_shape = new_slide.shapes.add_shape(
                                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
                            )
                            new_shape.text = shape.text_frame.text
                        elif hasattr(shape, "image"):
                            image_path = shape.image.filename
                            if os.path.exists(image_path):
                                new_slide.shapes.add_picture(image_path, shape.left, shape.top)
                            else:
                                print(f"Warning: Image file {image_path} not found. Skipping this image.")
            
            self.prs = new_prs
            print(f"{slide_name} deleted.")

    def create_presentation(self):
        """Save the PowerPoint presentation."""
        save_path, _ = QFileDialog.getSaveFileName(self, 'Save Presentation', '', 'PowerPoint files (*.pptx)')
        if save_path:
            self.prs.save(save_path)
            print(f"Presentation saved at {save_path}")


if __name__ == "__main__":
    app = QApplication(sys.argv)
    ui = ScriptureSlides()
    ui.show()
    sys.exit(app.exec_())
